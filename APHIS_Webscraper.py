#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Sun Apr 17 15:52:17 2022

@author: jackogozaly
"""

#Hello, if you've never used python before this should be fairly simple

#Firstly, make sure you have python installed. I reccomend downloading anaconda
#That way you'll have an enviroment to see the script in and won't have to bother
#As much with setting up packages. Download anaconda at the link below:
#https://docs.anaconda.com/anaconda/install/index.html


#This script relies on some packages, anaconda should come installed with most
#of these, but you'll have to download the following packages. Paste them in one 
#line at a time into the terminal and without the #

#pip install selenium
#pip install webdriver-manager
#pip install python-pptx


#This script utilizes the following packages:
#Used to create a browser window to navigate around in
from selenium import webdriver
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.common.by import By
from selenium.webdriver.chrome.service import Service
#Used for reading the website text
from bs4 import BeautifulSoup
import requests
#Used for general data manipulation and convenience
import time
import itertools
from tqdm import tqdm
import pandas as pd
import numpy as np
import re
from datetime import datetime
#Used to export the final result as a powerpoint deck
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt




# Welcome to Mr. Website Roboto, 

#Step 1: Identify some words that you think will be associated with 
#        websites that will need to be changed. 
keywords = ['PPQ 588', 'ePermit', 'PPQ 526', 'PPQ 525', 'PPQ 586', 'PPQ 587', 'PPQ 621', 
            'PPQ 585', 'PPQ 546', 'labels']

#Step 2: Identify some words that if present 100% mean the website needs to be 
#        changed. 
keywords2 = ['epermit', 'epermits']


#_____________________________________________________________________________


#Step 3: Run the script and let Mr. Website Roboto find what you're looking for!

#_____________________________________________________________________________

print("Beep boop, let me boot up google chrome and get to scraping!")
#Download our chrome driver
driver = webdriver.Chrome(service= Service(ChromeDriverManager().install()))

#Use the search engine on the sight to find relevant links
links_list = []
for i in range(len(keywords)):
    #Navigate to APHIS homepage
    driver.get("https://www.aphis.usda.gov/aphis/home/")
    #Set window size to maximum so we can find the searchbar
    driver.maximize_window()
    #Stop program for 3 seconds just to let stuff load
    time.sleep(1)
    
    #Search for our keywords
    search_bar = driver.find_element(by=By.CSS_SELECTOR, value= "#searchBox")
    search_bar.send_keys(keywords[i])
    search_bar.send_keys(Keys.RETURN)
    
    while True:
        html = driver.page_source
        soup = BeautifulSoup(html, "html.parser")
        #Create an empty list for our new links to go into
        links = []
        for link in soup.findAll('a'):
            links.append(link.get('href'))
        links_list.append(tuple(zip(links, [keywords[i] for x in range(len(links))])))
    
        try:
            #driver.find_element_by_class_name("next_page").click()
           driver.find_element(by=By.CLASS_NAME, value="next_page").click()
    
        except:
            break

#Exit out of our driver
driver.quit()



#Combine all the lists together
merged = list(itertools.chain(*links_list))
merged = [item for item in merged if item[0].startswith('https')]

#Find URLs that link to a pdf
pdf_links = [item for item in merged if item[0].endswith('pdf')]
#Find URLs that link to an excel sheet
excel_links = [item for item in merged if item[0].endswith('xlsx')]

#Filter out pdf and excel links, and create a list of clean urls to comb through
z = pdf_links + excel_links
web_links = [item for item in merged if item not in z]
links_to_search = [i[0] for i in web_links]
links_to_search = list(set(links_to_search))


#Convert keywords to search for to lowercase
for i in range(len(keywords2)):
    keywords2[i] = keywords2[i].lower()

#Create empty list for us to put the final links into
links_to_fix = []
offending_word = []
parent_elem = []
links = []
why_it_needs_revision = []

print("Beep boop, now to search the links for your keywords")
#Go to every website we previously got, and search for the keywords
for i in tqdm(range(len(links_to_search))):
  URL= links_to_search[i]
  page = requests.get(URL)
  soup = BeautifulSoup(page.content, 'html.parser')
  

  for link in soup.findAll('a'):
     links.append(link.get('href'))

  
  #Locate where the keywords are being said
  text = soup.get_text().lower().strip()
  pattern = re.compile(r'\b(?:%s)\b' % '|'.join(keywords2), re.IGNORECASE)
  searched_text = soup.find_all(text=pattern)
  
  
  if searched_text is not None: 
      for i in range(len(searched_text)):
          links_to_fix.append(URL)
          offending_word.append(searched_text[i])
          parent_elem.append(searched_text[i].parent.text)
          try:              
              if searched_text[0] == "Permits (ePermits and eFile)" and len(searched_text) == 1:
                  why_it_needs_revision.append("Banner says ePermit")
              else: 
                  why_it_needs_revision.append("Your keywords are mentioned somewhere on this page.")
          except: 
              why_it_needs_revision.append("I don't think ePermit is here, but you should verify")



#Deduplicate the links and filter out links we don't want to check for epermits
links = list(set(links))
links = pd.Series(links)
links = links[links.notnull()]
links = links[~links.str.contains('|'.join(['collapse', '#', 'mailto', '.pdf', '.csv',
                                         '.xlsx']))]

links = np.where(links.str.startswith('https'), 
                                links, 
                                'https://www.aphis.usda.gov/' + links)
links = [x for x in links if pd.isnull(x) == False]
links = list(links)
links = list(filter(lambda x: x.startswith("https://www.aphis.usda"), links))
links = [x for x in links if x not in links_to_search]

     
print("beep boop, oh man this will take a while. Let me look through all the USDA links I could find.")
for i in tqdm(range(len(links))):
  URL= links[i]
  page = requests.get(URL)
  soup = BeautifulSoup(page.content, 'html.parser')
    
  #Locate where the keywords are being said
  text = soup.get_text().lower().strip()
  pattern = re.compile(r'\b(?:%s)\b' % '|'.join(keywords2), re.IGNORECASE)
  searched_text = soup.find_all(text=pattern)
  
  
  if searched_text is not None: 
      for i in range(len(searched_text)):
          links_to_fix.append(URL)
          offending_word.append(searched_text[i])
          parent_elem.append(searched_text[i].parent.text)
          
          try:              
              if searched_text[0] == "Permits (ePermits and eFile)" and len(searched_text) == 1:
                  why_it_needs_revision.append("Banner says ePermit")
              else: 
                  why_it_needs_revision.append("Your keywords are mentioned somewhere on this page.")
          except: 
              why_it_needs_revision.append("I don't think ePermit is here, but you should verify")


df2 = pd.DataFrame({'URL': links_to_fix,
                    'Offending_Text': offending_word,
                    'Parent_Element': parent_elem,
                    'Why_it_needs_revision': why_it_needs_revision})


#Have it so we're only showing the banner change once
banner_rows = df2[df2['Why_it_needs_revision'] == 'Banner says ePermit']
non_banner_rows = df2[df2['Why_it_needs_revision'] != 'Banner says ePermit']

if banner_rows.shape[0] > 0: 
    banner_rows = banner_rows.iloc[0]
    all_links_to_change = pd.concat([banner_rows,non_banner_rows])
    
else:
    all_links_to_change = non_banner_rows

all_links_to_change = all_links_to_change.drop_duplicates()



#Function we will use to highlight the text we have
def highlight(element, effect_time, color, border):
    """Highlights (blinks) a Selenium Webdriver element"""
    driver = element._parent
    def apply_style(s):
        driver.execute_script("arguments[0].setAttribute('style', arguments[1]);",
                              element, s)
    apply_style("border: {0}px solid {1};".format(border, color))
    time.sleep(effect_time)


#Organize the final list of all the URLs we need to change
urls_to_screenshot = list(all_links_to_change['URL'].unique())
urls_to_screenshot = [x for x in urls_to_screenshot if pd.isnull(x) == False]

#_________________________Presentation Making Section__________________________

#Create the first slide of our presentation with all the necessary text
X = Presentation()
Layout = X.slide_layouts[0] 
first_slide = X.slides.add_slide(Layout)
first_slide.shapes.title.text = "Website Text Changes"
first_slide.placeholders[1].text = "Created by Mr. Website Roboto"

textbox = first_slide.shapes.add_textbox(Inches(2), Inches(5.21),Inches(6), Inches(2)) 
textframe = textbox.text_frame
textframe.word_wrap = True
paragraph = textframe.add_paragraph()
paragraph.text = f"On {datetime.today().strftime('%Y-%m-%d')}, I searched the APHIS website for all links related\
 to {str(keywords)[1:-1]}. In total, I examined {len(links) + len(links_to_search)} websites for your keywords and\
 I found {len(urls_to_screenshot)} websites that need to be altered. "
paragraph.font.size = Pt(12.5)


#Boot up google chrome to start taking screenshots
driver = webdriver.Chrome(service= Service(ChromeDriverManager().install()))
driver.set_window_size(1536, 1536)

for i in range(len(urls_to_screenshot)): 
    URL = urls_to_screenshot[i]
    driver.get(URL)
    time.sleep(2)
    screenshot = driver.save_screenshot('my_screenshot.png')
    
    df_subset = all_links_to_change[all_links_to_change['URL'] == URL]
    df_subset = df_subset.reset_index(drop=True)
    for i in range(df_subset.shape[0]):
        try: 
            open_window_elem = driver.find_elements_by_xpath(f"//*[contains(text(), '{df_subset['Offending_Text'][i]}')]")
            for i in range(len(open_window_elem)): 
                highlight(open_window_elem[i], 3, "orange", 5)
        except: 
            continue
    
    time.sleep(.5)
    
    screenshot2 = driver.save_screenshot('my_screenshot2.png')
    
    Second_Layout = X.slide_layouts[6]
    second_slide = X.slides.add_slide(Second_Layout)
    
    #Add in URL text
    textbox = second_slide.shapes.add_textbox(Inches(10), Inches(0),Inches(2.5), Inches(1)) 
    textframe = textbox.text_frame
    textframe.word_wrap = True
    paragraph = textframe.add_paragraph()
    paragraph.text = 'URL: ' + URL
    paragraph.font.size = Pt(12.5)
    
    textbox = second_slide.shapes.add_textbox(Inches(10), Inches(3.25),Inches(2.5), Inches(1)) 
    textframe = textbox.text_frame
    textframe.word_wrap = True
    paragraph = textframe.add_paragraph()
    paragraph.text = f'Why I flagged this: I counted {df_subset.shape[0]} elements that need to be changed here.'
    paragraph.font.size = Pt(12.5)
    
    textbox = second_slide.shapes.add_textbox(Inches(-6.1), Inches(0),Inches(2.5), Inches(1)) 
    textframe = textbox.text_frame
    textframe.word_wrap = True
    paragraph = textframe.add_paragraph()
    paragraph.text = 'Current Version'
    paragraph.font.size = Pt(12.5)
    
    
    pic = second_slide.shapes.add_picture('my_screenshot.png', left= Inches(-6.1), top= Inches(.75),
                                          width=Inches(6), height=Inches(6))
    
    pic2 = second_slide.shapes.add_picture('my_screenshot2.png', left= Inches(1.65), top= Inches(0),
                                          width=Inches(7.5), height=Inches(7.5))



X.save("APHIS_presentation.pptx")


driver.quit()


print("I'm all done!")
